import pathlib,options,os, subprocess
from drawer import pptx2png
from pptx import Presentation as pptxPresent
import logger_factory,logging

logging.getLogger(__file__)
class designer():
    def __init__(self ,data_pptxs):
        self.data_pptxs = data_pptxs

    def open_template_pptx(self,templateFileName):
        filePath = str(pathlib.Path(options.pptxs_path).joinpath(templateFileName))
        return pptxPresent(filePath)

    def fill_pptx_with_data(self,pptxObject,category,title,description,index=None,example=None,exampleOutput=None):

        logging.debug("filling in index slide -> Start")
        try:
            if index:
                for field in pptxObject.slides[0].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{index}':
                            field.text_frame.paragraphs[0].runs[0].text = f"{int(index):02d}"
                    except:
                        continue
                for field in pptxObject.slides[1].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{index}':
                            field.text_frame.paragraphs[0].runs[0].text = f"{int(index):02d}"
                    except:
                        continue
                for field in pptxObject.slides[2].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{index}':
                            field.text_frame.paragraphs[0].runs[0].text = f"{int(index):02d}"
                    except:
                        continue
        except Exception as msg:
            logging.error(f"Error in filling index slide -> {msg} \n{pptxObject} -> {title} ")
            return False
        else:
            logging.debug("filling in index slide -> OK")



        logging.debug("filling in category slide -> Start")
        try:
            if category:
                for field in pptxObject.slides[0].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{category}':
                            field.text_frame.paragraphs[0].runs[0].text = category
                    except:
                        continue
                for field in pptxObject.slides[1].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{category}':
                            field.text_frame.paragraphs[0].runs[0].text = category
                    except:
                        continue
                for field in pptxObject.slides[2].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{category}':
                            field.text_frame.paragraphs[0].runs[0].text = category
                    except:
                        continue
        except Exception as msg:
            logging.error(f"Error in filling category slide -> {msg} \n{pptxObject} -> {title} ")
            return False
        else:
            logging.debug("filling in category slide -> OK")

        logging.debug("filling in title slide -> Start")

        try:
            if title:
                for field in pptxObject.slides[0].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{title}':
                            field.text_frame.paragraphs[0].runs[0].text = title
                    except:
                        continue
                for field in pptxObject.slides[1].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{title}':
                            field.text_frame.paragraphs[0].runs[0].text = title
                    except:
                        continue
                for field in pptxObject.slides[2].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{title}':
                            field.text_frame.paragraphs[0].runs[0].text = title
                    except:
                        continue
        except Exception as msg:
            logging.error(f"Error in filling title slide -> {msg} \n{pptxObject} -> {title} ")
            return False
        else:
            logging.debug("filling in title slide -> OK")

        logging.debug("filling in description slide")
        try:
            if description:
                for field in pptxObject.slides[1].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{description}':
                            field.text_frame.paragraphs[0].runs[0].text = description
                    except:
                        continue
        except Exception as msg:
            logging.error(f"Error in filling description slide -> {msg} \n{pptxObject} -> {description} ")
            return False
        else:
            logging.debug("filling in description slide -> OK")

        logging.debug("filling in example slide")
        try:
            if example:
                for field in pptxObject.slides[2].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{example}':
                            field.text_frame.paragraphs[0].runs[0].text = example
                    except:
                        continue
        except Exception as msg:
            logging.error(f"Error in filling example slide -> {msg} \n{pptxObject} -> {example} ")
            return False
        else:
            logging.debug("filling in example slide -> OK")

        logging.debug("filling in example output slide")
        try:
            if exampleOutput:
                for field in pptxObject.slides[0].shapes:
                    try:
                        if field.text_frame.paragraphs[0].runs[0].text == '{exampleOutput}':
                            field.text_frame.paragraphs[0].runs[0].text = exampleOutput
                    except:
                        continue
        except Exception as msg:
            logging.error(f"Error in filling exampleOutput slide -> {msg} \n{pptxObject} -> {exampleOutput} ")
            return False
        else:
            logging.debug("filling in exampleOutput slide -> OK")

        if not pathlib.Path(options.main_path).joinpath('pptx_exports').joinpath(category+'-' + title).exists():
            os.mkdir(str(pathlib.Path(options.main_path).joinpath('pptx_exports').joinpath(category+'-' + title)))

        logging.debug(f"saving pptx file -> {title}")
        try:
            pptxObject.save(pathlib.Path(options.main_path).joinpath('pptx_exports').joinpath(category+'-' + title).joinpath(title+'.pptx'))
        except Exception as msg:
            logging.error(f"error in saving pptx file -> {title} \n -> {msg}")
        else:
            logging.debug(f"success in saving pptx -> {title}")
            return True

    def export_pptx_into_png(self,category,title):
        try:
            try:
                pptx2png.convert2png(str(pathlib.Path(options.main_path).joinpath('pptx_exports').joinpath(category+'-' + title).joinpath(title+'.pptx')))
            except Exception as msg:
                logging.error("Couldn't run pptx2png, running by CMD")
                print(r"Couldn't run pptx2png, running by CMD")
                output = subprocess.call(['python',
                                          str(pathlib.Path(options.main_path).joinpath('drawer').joinpath('ppt2png_cmd.py')),
                                          str(pathlib.Path(options.main_path).joinpath('pptx_exports').joinpath(category+'-' + title).joinpath(title+'.pptx'))])
                print(output)
        except Exception as msg:
            logging.error(f"Error in exporting png from {pathlib.Path(options.main_path).joinpath('pptx_exports').joinpath(category+'-' + title).joinpath(title+'.pptx')}\n-> {msg}")
            return False
        else:
            return True

    def run_designer_for(self ,themefile, category,title ,description,index=None,example=None,expoutput=None):
        """
        steps : 1. open template pptx
                2. fill the pptx & save it in output folder
                3. run pptx2png in output folder and get png output
        :param themefile:
        :param category:
        :param title:
        :param description:
        :return:
        """
        thepptx = self.open_template_pptx(themefile)
        self.fill_pptx_with_data(thepptx,category,title,description,index,example,exampleOutput=expoutput)
        self.export_pptx_into_png(category,title)
        return True


    def design_data(self,index=True):
        logging.debug("starting designing data")
        for key,value in self.data_pptxs.items():
            print(1)
            for item in value:
                print(item)
                try:
                    #(['0', 'append()', 'Adds an element at   the end of the list'], 'test.pptx')
                    logging.debug(f"starting to export -> {item}")
                    if index:
                        try:
                            example = item[0][3]
                        except:
                            example = None
                        try:
                            exampleoutput = item[0][4]
                        except:
                            exampleoutput = None
                        self.run_designer_for(item[1],key.split('.')[0],item[0][1],item[0][2],item[0][0],example,exampleoutput)
                    else:
                        self.run_designer_for(item[1], key.split('.')[0], item[0][0], item[0][1])
                    logging.debug(f"exported -> {item}")
                except Exception as msg:
                    logging.debug(f"error in exporting -> {item}\n err -> {msg}")
                    return False
        return True